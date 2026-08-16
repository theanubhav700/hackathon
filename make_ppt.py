from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
RED    = RGBColor(0xCC, 0x00, 0x00)
GREY   = RGBColor(0x44, 0x44, 0x44)
LGREY  = RGBColor(0xF2, 0xF2, 0xF2)
DKGREY = RGBColor(0x22, 0x22, 0x22)

blank_layout = prs.slide_layouts[6]  # completely blank

# ─────────────────────────────────────────────────────────────
# Helper: add a filled rectangle
# ─────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape

# ─────────────────────────────────────────────────────────────
# Helper: add a text box
# ─────────────────────────────────────────────────────────────
def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return txb

# ─────────────────────────────────────────────────────────────
# Helper: add multi-para text box
# ─────────────────────────────────────────────────────────────
def add_multiline(slide, lines, l, t, w, h,
                  size=16, color=BLACK, align=PP_ALIGN.LEFT,
                  bold_first=False):
    """lines = list of (text, bold) tuples"""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, (txt, bld) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = txt
        run.font.size  = Pt(size)
        run.font.bold  = bld
        run.font.color.rgb = color
        run.font.name  = "Calibri"

# =============================================================
# SLIDE 1 — Title
# =============================================================
slide = prs.slides.add_slide(blank_layout)

# dark background
add_rect(slide, 0, 0, 13.33, 7.5, DARK)

# red top bar
add_rect(slide, 0, 0, 13.33, 0.12, RED)
# red bottom bar
add_rect(slide, 0, 7.38, 13.33, 0.12, RED)

# Title
add_text(slide, "ResQ", 1, 1.8, 11.33, 1.4,
         size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text(slide, "Smart Ambulance Route Optimization & Hospital Pre-Admission",
         1, 3.4, 11.33, 1.0,
         size=24, bold=False, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

# Tagline
add_text(slide, "Smart India Hackathon 2026",
         1, 4.6, 11.33, 0.6,
         size=18, bold=False, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)

# Slide number
add_text(slide, "1 / 8", 12.5, 7.1, 0.7, 0.3, size=10,
         color=RGBColor(0x66,0x66,0x66), align=PP_ALIGN.RIGHT)


# =============================================================
# SLIDE 2 — The Problem
# =============================================================
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 0.08, RED)
add_rect(slide, 0, 7.42, 13.33, 0.08, RED)

add_text(slide, "The Problem", 0.5, 0.25, 12, 0.7,
         size=32, bold=True, color=DARK, align=PP_ALIGN.LEFT)
add_rect(slide, 0.5, 0.95, 2.5, 0.04, RED)

add_multiline(slide, [
    ("Every minute matters in a medical emergency.", True),
    ("", False),
    ("But today, two critical problems make ambulances slow and hospitals unprepared:", False),
], 0.5, 1.2, 12.3, 1.2, size=18, color=DKGREY)

# Box 1
add_rect(slide, 0.5, 2.6, 5.8, 2.8, LGREY)
add_multiline(slide, [
    ("Problem 1", True),
    ("", False),
    ("Ambulances get stuck in traffic.", False),
    ("", False),
    ("There is no system to automatically clear the road for an ambulance. Traffic signals work on normal cycles and do not know an ambulance is approaching.", False),
], 0.7, 2.75, 5.4, 2.5, size=15, color=DKGREY)

# Box 2
add_rect(slide, 6.9, 2.6, 5.9, 2.8, LGREY)
add_multiline(slide, [
    ("Problem 2", True),
    ("", False),
    ("Hospitals are unprepared for critical patients.", False),
    ("", False),
    ("The ER team has no information about the patient's condition before the ambulance arrives. Precious minutes are wasted at the hospital door.", False),
], 7.1, 2.75, 5.5, 2.5, size=15, color=DKGREY)

add_text(slide, "Result: Patients lose golden minutes — the most critical window for survival.",
         0.5, 5.65, 12.3, 0.6,
         size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)

add_text(slide, "2 / 8", 12.5, 7.1, 0.7, 0.3, size=10,
         color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.RIGHT)


# =============================================================
# SLIDE 3 — Our Solution (Overview)
# =============================================================
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 0.08, RED)
add_rect(slide, 0, 7.42, 13.33, 0.08, RED)

add_text(slide, "Our Solution — ResQ", 0.5, 0.25, 12, 0.7,
         size=32, bold=True, color=DARK)
add_rect(slide, 0.5, 0.95, 3.0, 0.04, RED)

add_text(slide,
         "ResQ connects the patient, ambulance driver, admin dispatcher, and hospital ER into one real-time platform — from the moment a patient needs help to the moment they reach the ER.",
         0.5, 1.15, 12.3, 0.9, size=17, color=DKGREY)

# 4 boxes
boxes = [
    ("Patient", "Books ambulance in seconds with one tap. Location is detected automatically."),
    ("Admin", "Sees all active emergencies live and manages the entire ambulance fleet."),
    ("Driver", "Gets the booking instantly, navigates to the patient, and sends patient data to the hospital while driving."),
    ("Hospital ER", "Receives live patient vitals and ECG before the ambulance even arrives."),
]

positions = [(0.4, 2.3), (3.55, 2.3), (6.7, 2.3), (9.85, 2.3)]
for (title, desc), (lx, ty) in zip(boxes, positions):
    add_rect(slide, lx, ty, 3.0, 3.6, LGREY)
    add_multiline(slide, [
        (title, True),
        ("", False),
        (desc, False),
    ], lx+0.15, ty+0.15, 2.7, 3.3, size=14, color=DKGREY)

add_text(slide, "3 / 8", 12.5, 7.1, 0.7, 0.3, size=10,
         color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.RIGHT)


# =============================================================
# SLIDE 4 — How Patient Books (Feature 1)
# =============================================================
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 0.08, RED)
add_rect(slide, 0, 7.42, 13.33, 0.08, RED)

add_text(slide, "Feature 1 — Patient Booking", 0.5, 0.25, 12, 0.7,
         size=32, bold=True, color=DARK)
add_rect(slide, 0.5, 0.95, 3.5, 0.04, RED)

steps = [
    ("Step 1", "Patient opens ResQ and their location is detected automatically using GPS."),
    ("Step 2", "Patient selects the type of emergency — Cardiac Arrest, Accident, Stroke, Maternity, and more."),
    ("Step 3", "Available ambulances near the patient are shown with distance and estimated arrival time."),
    ("Step 4", "Patient confirms the booking. The driver receives an alert instantly."),
    ("Step 5", "Patient sees live status — driver notified, driver accepted, driver on the way."),
]

for i, (step, desc) in enumerate(steps):
    ty = 1.2 + i * 1.1
    add_rect(slide, 0.5, ty, 1.2, 0.75, DARK)
    add_text(slide, step, 0.5, ty, 1.2, 0.75,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_multiline(slide, [(desc, False)],
                  1.9, ty+0.05, 10.8, 0.7, size=15, color=DKGREY)

add_text(slide, "4 / 8", 12.5, 7.1, 0.7, 0.3, size=10,
         color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.RIGHT)


# =============================================================
# SLIDE 5 — Green Corridor (Feature 2)
# =============================================================
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 0.08, RED)
add_rect(slide, 0, 7.42, 13.33, 0.08, RED)

add_text(slide, "Feature 2 — Green Corridor", 0.5, 0.25, 12, 0.7,
         size=32, bold=True, color=DARK)
add_rect(slide, 0.5, 0.95, 3.5, 0.04, RED)

add_text(slide, "Solves Problem 1 — Ambulance stuck in traffic",
         0.5, 1.1, 12, 0.45, size=17, bold=True, color=RED)

add_text(slide,
         "When the driver activates Green Corridor, the system places traffic signals along the entire route to the hospital. As the ambulance moves forward, signals automatically turn green ahead of it and reset to red behind it. The ambulance gets a clear road all the way.",
         0.5, 1.65, 12.3, 1.1, size=16, color=DKGREY)

add_text(slide, "How it works:", 0.5, 2.85, 12, 0.45, size=16, bold=True, color=DARK)

how = [
    "The route from the ambulance to the hospital is calculated using real road data.",
    "7 virtual traffic signals are placed along this route.",
    "When the ambulance is 800 metres away from a signal — it turns yellow.",
    "When the ambulance is 400 metres away — it turns green.",
    "Once the ambulance passes — the signal resets to red.",
    "All of this happens automatically with no manual control needed.",
]
for i, line in enumerate(how):
    ty = 3.3 + i * 0.6
    add_rect(slide, 0.5, ty+0.12, 0.22, 0.22, RED)
    add_text(slide, line, 0.9, ty, 11.9, 0.55, size=14, color=DKGREY)

add_text(slide, "5 / 8", 12.5, 7.1, 0.7, 0.3, size=10,
         color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.RIGHT)


# =============================================================
# SLIDE 6 — Live Telemetry & Hospital ER (Feature 3)
# =============================================================
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 0.08, RED)
add_rect(slide, 0, 7.42, 13.33, 0.08, RED)

add_text(slide, "Feature 3 — Live Telemetry to Hospital ER", 0.5, 0.25, 12, 0.7,
         size=30, bold=True, color=DARK)
add_rect(slide, 0.5, 0.95, 4.5, 0.04, RED)

add_text(slide, "Solves Problem 2 — Hospital unprepared for incoming patient",
         0.5, 1.1, 12, 0.45, size=17, bold=True, color=RED)

add_text(slide,
         "While the ambulance is on the way, the driver enters the patient's medical vitals. This data is sent live to the hospital ER. The ER team sees everything in real time — no waiting until arrival.",
         0.5, 1.65, 12.3, 0.9, size=16, color=DKGREY)

# Left column — Driver sends
add_rect(slide, 0.5, 2.7, 5.8, 3.6, LGREY)
add_text(slide, "Driver Sends (from ambulance)", 0.65, 2.82, 5.5, 0.5,
         size=15, bold=True, color=DARK)
driver_sends = [
    "Heart Rate",
    "Blood Oxygen Level (SpO2)",
    "Blood Pressure",
    "Body Temperature",
    "Patient Condition — Critical / Serious / Stable / Normal",
    "Live ECG waveform streamed continuously",
    "Pre-arrival alert with patient summary before reaching hospital",
]
for i, item in enumerate(driver_sends):
    add_text(slide, "- " + item, 0.75, 3.3 + i*0.42, 5.3, 0.4, size=13, color=DKGREY)

# Right column — Hospital receives
add_rect(slide, 7.0, 2.7, 5.8, 3.6, LGREY)
add_text(slide, "Hospital ER Receives (live)", 7.15, 2.82, 5.5, 0.5,
         size=15, bold=True, color=DARK)
hospital_gets = [
    "All vitals update in real time on screen",
    "ECG waveform displayed on a live monitor",
    "Critical alerts highlighted automatically",
    "Pre-alert notification before ambulance arrives",
    "ER team can prepare — Trauma Bay, Blood Bank, ICU",
    "No login needed — just enter the booking ID",
]
for i, item in enumerate(hospital_gets):
    add_text(slide, "- " + item, 7.1, 3.3 + i*0.42, 5.5, 0.4, size=13, color=DKGREY)

add_text(slide, "6 / 8", 12.5, 7.1, 0.7, 0.3, size=10,
         color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.RIGHT)


# =============================================================
# SLIDE 7 — All Features Summary
# =============================================================
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 0.08, RED)
add_rect(slide, 0, 7.42, 13.33, 0.08, RED)

add_text(slide, "What ResQ Can Do — Full Feature List", 0.5, 0.25, 12, 0.7,
         size=30, bold=True, color=DARK)
add_rect(slide, 0.5, 0.95, 5.0, 0.04, RED)

col1 = [
    ("Patient Side", True),
    ("Auto GPS location detection", False),
    ("8 emergency types to choose from", False),
    ("See nearby ambulances with ETA", False),
    ("Live booking status updates", False),
    ("Full booking history", False),
    ("", False),
    ("Admin Side", True),
    ("Live view of all active emergencies", False),
    ("Manage entire ambulance fleet", False),
    ("Add and remove drivers", False),
    ("Analytics and activity logs", False),
]

col2 = [
    ("Driver Side", True),
    ("Receive emergency requests live", False),
    ("Accept or reject bookings", False),
    ("GPS navigation to patient", False),
    ("Green Corridor — auto traffic clearance", False),
    ("Enter and send patient vitals to ER", False),
    ("Stream live ECG to hospital", False),
    ("Send pre-arrival alert to hospital", False),
    ("Download patient vitals as PDF report", False),
    ("", False),
    ("Hospital Side", True),
    ("Live vitals and ECG monitor", False),
    ("Pre-alert inbox before patient arrives", False),
    ("One-click ER preparation actions", False),
]

add_multiline(slide, col1, 0.5, 1.15, 6.0, 6.0, size=14, color=DKGREY)
add_multiline(slide, col2, 6.9, 1.15, 6.0, 6.0, size=14, color=DKGREY)

add_text(slide, "7 / 8", 12.5, 7.1, 0.7, 0.3, size=10,
         color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.RIGHT)


# =============================================================
# SLIDE 8 — Problem vs Solution Summary
# =============================================================
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
add_rect(slide, 0, 0, 13.33, 0.08, RED)
add_rect(slide, 0, 7.42, 13.33, 0.08, RED)

add_text(slide, "Problem vs Solution", 0.5, 0.25, 12, 0.7,
         size=32, bold=True, color=DARK)
add_rect(slide, 0.5, 0.95, 3.0, 0.04, RED)

# Table headers
add_rect(slide, 0.5,  1.15, 5.5, 0.55, DARK)
add_rect(slide, 6.2,  1.15, 6.6, 0.55, DARK)
add_text(slide, "Problem", 0.5, 1.15, 5.5, 0.55,
         size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "How ResQ Solves It", 6.2, 1.15, 6.6, 0.55,
         size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ("Ambulance stuck in traffic — no one clears the road",
     "Green Corridor automatically turns signals green 400m ahead of the ambulance along the entire route"),
    ("Hospital has no information before patient arrives",
     "Driver sends live vitals and ECG from the ambulance. Hospital sees everything in real time before arrival"),
    ("No system to book an ambulance quickly",
     "Patient books in seconds — GPS detects location automatically, nearby ambulances shown instantly"),
    ("Driver has no navigation or patient data tools",
     "Driver gets turn-by-turn navigation, can record vitals, stream ECG, send pre-alert, and download PDF report"),
    ("Admin cannot see what is happening in the field",
     "Admin dashboard shows all active emergencies live, full fleet status, and driver locations in real time"),
]

for i, (prob, sol) in enumerate(rows):
    ty = 1.8 + i * 1.05
    bg = LGREY if i % 2 == 0 else WHITE
    add_rect(slide, 0.5, ty, 5.5, 0.95, bg)
    add_rect(slide, 6.2, ty, 6.6, 0.95, bg)
    add_text(slide, prob, 0.6, ty+0.05, 5.3, 0.88, size=13, color=DKGREY)
    add_text(slide, sol,  6.3, ty+0.05, 6.4, 0.88, size=13, color=DKGREY)

add_text(slide, "8 / 8", 12.5, 7.1, 0.7, 0.3, size=10,
         color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.RIGHT)


# =============================================================
# Save
# =============================================================
out = r"c:\Users\Anubhav Tiwari\Downloads\hackathon\ResQ_Presentation.pptx"
prs.save(out)
print("Saved:", out)
